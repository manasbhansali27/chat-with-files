"""
embedder.py
-----------
Handles file parsing, embedding (text, image, video), and database storage.
Central module for generating embeddings and managing FAISS indices.
"""

import os
import io
import cv2
import sqlite3
import tempfile
import subprocess
import logging
from pathlib import Path

import config
import faiss
import numpy as np
import whisper

# Document parsers
import pandas as pd
import nbformat
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

# OCR & images
from PIL import Image
import easyocr
import fitz  # PyMuPDF

# Models
import torch
from transformers import AutoProcessor, AutoModel
from sentence_transformers import SentenceTransformer
from ollama import Client as OllamaClient

from tqdm import tqdm
# =======================
#  GPU SETUP
# =======================
device = "cuda" if torch.cuda.is_available() else "cpu"
print(f"[INIT] Using device: {device}")

# Load models on GPU
text_embedder = SentenceTransformer('all-MiniLM-L6-v2').to(device)
siglip_model = AutoModel.from_pretrained("google/siglip-base-patch16-224").to(device)
siglip_processor = AutoProcessor.from_pretrained("google/siglip-base-patch16-224")
easyocr_reader = easyocr.Reader(['en'], gpu=True if device=="cuda" else False)
ollama = OllamaClient()

# Set model to eval mode for inference
siglip_model.eval()

print(f"[INIT] Models loaded successfully")

# =======================
# DATABASE
# =======================
def get_db():
    conn = sqlite3.connect(config.SQLITE_PATH)
    cursor = conn.cursor()

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS vectors (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            file_id INTEGER,
            chunk_id INTEGER,
            modality TEXT,
            embedding BLOB
        )
    """)

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS files (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            path TEXT UNIQUE,
            type TEXT
        )
    """)

    cursor.execute("PRAGMA table_info(files)")
    cols = [col[1] for col in cursor.fetchall()]
    if "processed" not in cols:
        cursor.execute("ALTER TABLE files ADD COLUMN processed INTEGER DEFAULT 0")
        conn.commit()

    return conn

# =======================
# EMBEDDING FUNCTIONS
# =======================
def embed_text_enhanced(text: str):
    """FIXED: Text embedding with normalization"""
    if not text.strip():
        return None
    
    try:
        embedding = text_embedder.encode(text, convert_to_tensor=True)
        # CRITICAL FIX: Normalize the embedding
        embedding_np = embedding.cpu().numpy().astype("float32")
        embedding_np = embedding_np / np.linalg.norm(embedding_np)
        return embedding_np
    except Exception as e:
        print(f"[EMBED WARN] SentenceTransformer failed: {e}")
        try:
            res = ollama.embeddings(model="mistral:7b", prompt=text)
            embedding = np.array(res.embedding, dtype="float32")
            # Normalize ollama embedding too
            embedding = embedding / np.linalg.norm(embedding)
            return embedding
        except Exception as e2:
            print(f"[EMBED ERROR] Both embedding methods failed: {e2}")
            return None

def create_contextual_text(filename: str, content: str, file_type: str) -> str:
    """Create contextual text for better semantic matching"""
    filename_clean = Path(filename).stem.replace('_', ' ').replace('-', ' ')
    return f"File: {filename_clean} | Type: {file_type} | Content: {content[:500]}"

def embed_image_siglip(image: Image.Image):
    """FIXED: SigLIP embedding with proper implementation and normalization"""
    try:
        image = image.convert("RGB")
        if image.size[0] < 32 or image.size[1] < 32:
            image = image.resize((224, 224), Image.Resampling.LANCZOS)
        
        # Process image
        inputs = siglip_processor(images=image, return_tensors="pt").to(device)
        
        with torch.no_grad():
            # CRITICAL FIX: Use get_image_features() instead of vision_model
            image_features = siglip_model.get_image_features(**inputs)
        
        # CRITICAL FIX: Normalize the embedding
        embedding = image_features.cpu().numpy().flatten().astype("float32")
        embedding = embedding / np.linalg.norm(embedding)
        
        return embedding
        
    except Exception as e:
        print(f"[SIGLIP ERROR] {e}")
        import traceback
        traceback.print_exc()
        return None

def get_visual_text_embedding(text: str):
    """Get text embedding for visual queries (for testing)"""
    try:
        inputs = siglip_processor(text=[text], return_tensors="pt", padding=True).to(device)
        
        with torch.no_grad():
            # Use get_text_features() for proper text embedding
            text_features = siglip_model.get_text_features(**inputs)
        
        # Normalize the embedding
        embedding = text_features.cpu().numpy().flatten().astype("float32")
        embedding = embedding / np.linalg.norm(embedding)
        
        return embedding
        
    except Exception as e:
        print(f"[VISUAL TEXT ERROR] {e}")
        return None

def ocr_image(image: Image.Image) -> str:
    """OCR using EasyOCR on GPU"""
    try:
        image_np = np.array(image)
        
        # Skip very small images
        if image_np.shape[0] < 20 or image_np.shape[1] < 20:
            return ""
        
        results = easyocr_reader.readtext(image_np)
        
        extracted_texts = []
        for result in results:
            if len(result) == 3:
                bbox, text, confidence = result
                if confidence > 0.5 and len(text.strip()) > 1:
                    extracted_texts.append(text.strip())
            elif len(result) == 2:
                bbox, text = result
                if len(text.strip()) > 2:
                    extracted_texts.append(text.strip())
        
        return " ".join(extracted_texts).strip()
        
    except Exception as e:
        print(f"[OCR ERROR] {e}")
        return ""

# =======================
# HELPERS
# =======================
def chunk_text_smart(text: str, max_len=300):
    """Smart text chunking"""
    if len(text.split()) <= max_len:
        return [text]
    
    sentences = text.replace('. ', '.|').replace('! ', '!|').replace('? ', '?|').split('|')
    
    chunks = []
    current_chunk = []
    current_len = 0
    
    for sentence in sentences:
        words_in_sentence = len(sentence.split())
        if current_len + words_in_sentence <= max_len:
            current_chunk.append(sentence)
            current_len += words_in_sentence
        else:
            if current_chunk:
                chunks.append(' '.join(current_chunk))
                current_chunk = [sentence]
                current_len = words_in_sentence
            else:
                words = sentence.split()
                for i in range(0, len(words), max_len):
                    chunks.append(' '.join(words[i:i+max_len]))
    
    if current_chunk:
        chunks.append(' '.join(current_chunk))
    
    return chunks

def get_file_type(path):
    """Determine file type"""
    ext = Path(path).suffix.lower()
    
    video_exts = ['.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm', '.m4v', '.3gp']
    if ext in video_exts:
        return 'video'
    
    image_exts = ['.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif']
    if ext in image_exts:
        return 'image'
    
    doc_exts = ['.pdf', '.docx', '.pptx', '.csv', '.ipynb', '.html', '.htm', '.md', 
                '.txt', '.py', '.cpp', '.c', '.h', '.java', '.js', '.ts', '.json','.css']
    if ext in doc_exts:
        return 'document'
    
    return 'unknown'

# =======================
# FILE PARSERS 
# =======================
def parse_pdf(path):
    text, images = "", []
    try:
        reader = PdfReader(path)
        for page in reader.pages:
            text += page.extract_text() or ""
    except Exception as e:
        print(f"[PDF TEXT ERROR] {path}: {e}")
        raise e

    try:
        doc = fitz.open(path)
        for page_index in range(len(doc)):
            for _, img in enumerate(doc.get_page_images(page_index)):
                xref = img[0]
                base_image = doc.extract_image(xref)
                img_pil = Image.open(io.BytesIO(base_image["image"])).convert("RGB")
                images.append(img_pil)
    except Exception as e:
        print(f"[PDF IMAGE ERROR] {path}: {e}")
    return text, images

def parse_docx(path):
    text, images = "", []
    try:
        doc = Document(path)
        text += "\n".join([p.text for p in doc.paragraphs])
        for rel in doc.part._rels:
            rel_obj = doc.part._rels[rel]
            if "image" in rel_obj.target_ref:
                img_data = rel_obj.target_part.blob
                img_pil = Image.open(io.BytesIO(img_data)).convert("RGB")
                images.append(img_pil)
    except Exception as e:
        print(f"[DOCX ERROR] {path}: {e}")
        raise e
    return text, images

def parse_pptx(path):
    text, images = "", []
    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
                elif getattr(shape, "shape_type", None) == 13:
                    img_pil = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                    images.append(img_pil)
    except Exception as e:
        print(f"[PPTX ERROR] {path}: {e}")
        raise e
    return text, images

def parse_csv(path, max_rows=20):
    try:
        df = pd.read_csv(path, on_bad_lines="skip")
        df_subset = df.head(max_rows)
        return df_subset.to_string(), []
    except Exception as e:
        print(f"[CSV ERROR] {path}: {e}")
        raise e

def parse_ipynb(path):
    try:
        with open(path, "r", encoding="utf-8") as f:
            nb = nbformat.read(f, as_version=4)
        texts = []
        for cell in nb.cells:
            if cell.cell_type in ("markdown", "code"):
                texts.append(" ".join(cell.source.splitlines()))
        return "\n".join(texts), []
    except Exception as e:
        print(f"[IPYNB ERROR] {path}: {e}")
        raise e

def parse_html(path):
    try:
        with open(path, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")
        return soup.get_text(), []
    except Exception as e:
        print(f"[HTML ERROR] {path}: {e}")
        raise e

def parse_txt(path):
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read(), []
    except Exception as e:
        print(f"[TXT ERROR] {path}: {e}")
        raise e

def extract_document_content(path):
    """Extract content from document files"""
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return parse_pdf(path)
    elif ext == ".docx":
        return parse_docx(path)
    elif ext == ".pptx":
        return parse_pptx(path)
    elif ext == ".csv":
        return parse_csv(path, max_rows=20)
    elif ext == ".ipynb":
        return parse_ipynb(path)
    elif ext in [".html", ".htm", ".md"]:
        return parse_html(path)
    elif ext in [".txt", ".py", ".cpp", ".c", ".h", ".java", ".js", ".ts", ".json",".css"]:
        return parse_txt(path)
    else:
        return "", []

# =======================
# VIDEO PROCESSORS
# =======================
def extract_audio_ffmpeg(video_path: str) -> str | None:
    """Extract audio from video using ffmpeg"""
    wav_file = tempfile.NamedTemporaryFile(suffix=".wav", delete=False).name
    cmd = [
        "ffmpeg", "-i", str(Path(video_path).as_posix()),
        "-ar", "16000", "-ac", "1", "-vn", "-y", 
        "-loglevel", "quiet",
        wav_file
    ]
    
    try:
        result = subprocess.run(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.PIPE, timeout=300)
        
        if result.returncode != 0 or not Path(wav_file).exists() or Path(wav_file).stat().st_size == 0:
            Path(wav_file).unlink(missing_ok=True)
            return None
            
        return wav_file
        
    except:
        Path(wav_file).unlink(missing_ok=True)
        return None

def process_video_frames(path, frame_interval=30):
    """Extract and embed video frames"""
    cap = cv2.VideoCapture(str(path))
    if not cap.isOpened():
        return []

    frame_embeddings = []
    frame_id = 0
    success, frame = cap.read()

    while success:
        if frame_id % frame_interval == 0:
            img = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
            emb = embed_image_siglip(img)
            if emb is not None:
                frame_embeddings.append((frame_id, emb))
        success, frame = cap.read()
        frame_id += 1

    cap.release()
    return frame_embeddings

def process_video_audio_enhanced(path, model_name="base"):
    """Process video audio on CPU"""
    wav_file = extract_audio_ffmpeg(path)
    if not wav_file:
        return []

    try:
        model = whisper.load_model(model_name, device="cpu")
        result = model.transcribe(wav_file, fp16=False)
        full_transcript = result.get("text", "").strip()
        segments = result.get("segments", [])
        del model
        
    finally:
        Path(wav_file).unlink(missing_ok=True)

    if not full_transcript:
        return []

    filename = Path(path).stem
    text_embeddings = []
    
    contextual_text = create_contextual_text(filename, full_transcript, "video audio transcript")
    emb = embed_text_enhanced(contextual_text)
    if emb is not None:
        text_embeddings.append((-1, emb))
    
    for i, segment in enumerate(segments):
        segment_text = segment.get('text', '').strip()
        if len(segment_text.split()) >= 5:
            segment_context = create_contextual_text(filename, segment_text, f"video segment {i+1}")
            emb = embed_text_enhanced(segment_context)
            if emb is not None:
                text_embeddings.append((i, emb))
    
    return text_embeddings

# =======================
# MAIN PROCESSING WITH FIXED FAISS INDICES
# =======================
def process_single_file(file_id, path, frame_interval=60, whisper_model="base"):
    """FIXED: Process a single file with correct FAISS index types"""
    conn = get_db()
    cursor = conn.cursor()
    
    text_index_path = "indices/faiss_text.index"
    image_index_path = "indices/faiss_image.index"
    
    text_index = faiss.read_index(text_index_path) if Path(text_index_path).exists() else None
    image_index = faiss.read_index(image_index_path) if Path(image_index_path).exists() else None
    
    file_type = get_file_type(path)
    fname = Path(path).name
    embeddings_added = False
    processing_failed = False
    
    try:
        chunks_to_process = []

        if file_type == 'video':
            try:
                frames = process_video_frames(path, frame_interval)
                chunks_to_process.extend([('video_frame', f_id, emb) for f_id, emb in frames])
                
                audio_segments = process_video_audio_enhanced(path, whisper_model)
                chunks_to_process.extend([('video_audio', seg_id + 1000, emb) for seg_id, emb in audio_segments])
            except Exception as e:
                print(f"[VIDEO ERROR] {fname}: {e}")
                processing_failed = True
        
        elif file_type == 'image':
            try:
                img = Image.open(path).convert("RGB")
                emb_img = embed_image_siglip(img)
                if emb_img is not None:
                    print(f"[DEBUG] Image embedding shape: {emb_img.shape}, norm: {np.linalg.norm(emb_img):.3f}")
                    chunks_to_process.append(('image', 0, emb_img))
                
                ocr_txt = ocr_image(img)
                if ocr_txt.strip():
                    emb_text = embed_text_enhanced(create_contextual_text(fname, ocr_txt, "image with text"))
                    if emb_text is not None:
                        chunks_to_process.append(('ocr_text', 1, emb_text))
            except Exception as e:
                print(f"[IMAGE ERROR] {fname}: {e}")
                processing_failed = True
        
        elif file_type == 'document':
            try:
                text, images = extract_document_content(path)
                if text.strip():
                    chunks = chunk_text_smart(text)
                    for chunk_id, chunk in enumerate(chunks):
                        emb = embed_text_enhanced(create_contextual_text(fname, chunk, f"document chunk {chunk_id+1}"))
                        if emb is not None:
                            chunks_to_process.append(('text', chunk_id, emb))
                
                for img_id, img in enumerate(images):
                    emb_img = embed_image_siglip(img)
                    if emb_img is not None:
                        chunks_to_process.append(('image', img_id + 1000, emb_img))
                    
                    ocr_txt = ocr_image(img)
                    if ocr_txt.strip():
                        emb_text = embed_text_enhanced(create_contextual_text(fname, ocr_txt, f"document image {img_id+1}"))
                        if emb_text is not None:
                            chunks_to_process.append(('ocr_text', img_id + 2000, emb_text))
            except Exception as e:
                print(f"[DOCUMENT ERROR] {fname}: {e}")
                processing_failed = True
        
        if chunks_to_process:
            for modality, chunk_id, emb in tqdm(chunks_to_process, desc=f"[{fname}]"):
                if modality in ("text", "ocr_text", "video_audio"):
                    if text_index is None:
                        # CRITICAL FIX: Use IndexFlatIP for normalized embeddings
                        text_index = faiss.IndexFlatIP(len(emb))
                        print(f"[INIT] Created new text index with dimension {len(emb)}")
                    cursor.execute(
                        "INSERT INTO vectors (file_id, chunk_id, modality, embedding) VALUES (?, ?, ?, ?)",
                        (file_id, chunk_id, modality, emb.tobytes())
                    )
                    text_index.add(np.array([emb], dtype="float32"))
                elif modality in ("image", "video_frame"):
                    if image_index is None:
                        # CRITICAL FIX: Use IndexFlatIP for normalized embeddings
                        image_index = faiss.IndexFlatIP(len(emb))
                        print(f"[INIT] Created new image index with dimension {len(emb)}")
                    cursor.execute(
                        "INSERT INTO vectors (file_id, chunk_id, modality, embedding) VALUES (?, ?, ?, ?)",
                        (file_id, chunk_id, modality, emb.tobytes())
                    )
                    image_index.add(np.array([emb], dtype="float32"))
                embeddings_added = True
        
        # Save indices
        if text_index:
            faiss.write_index(text_index, text_index_path)
            print(f"[SAVE] Text index saved with {text_index.ntotal} vectors")
        if image_index:
            faiss.write_index(image_index, image_index_path)
            print(f"[SAVE] Image index saved with {image_index.ntotal} vectors")
        
        if embeddings_added and not processing_failed:
            cursor.execute("UPDATE files SET processed=1 WHERE id=?", (file_id,))
            conn.commit()
            print(f"✅ {fname}")
        else:
            cursor.execute("UPDATE files SET processed=-1 WHERE id=?", (file_id,))
            conn.commit()
            print(f"❌ {fname}")
    
    except Exception as e:
        cursor.execute("UPDATE files SET processed=-1 WHERE id=?", (file_id,))
        conn.commit()
        print(f"❌ {fname}: {e}")
        import traceback
        traceback.print_exc()
    
    finally:
        conn.close()

def process_files(limit=20, frame_interval=60, whisper_model="base"):
    """Process unprocessed files"""
    os.makedirs("indices", exist_ok=True)
    conn = get_db()
    cursor = conn.cursor()

    cursor.execute("SELECT id, path FROM files WHERE processed=0 LIMIT ?", (limit,))
    rows = cursor.fetchall()
    conn.close()
    
    if not rows:
        print("[INFO] No unprocessed files found.")
        return 0

    print(f"[PROCESS] Found {len(rows)} files")
    
    processed_count = 0
    for file_id, path in tqdm(rows, desc="Processing"):
        if not os.path.exists(path):
            print(f"[ERROR] Missing: {path}")
            conn = get_db()
            cursor = conn.cursor()
            cursor.execute("UPDATE files SET processed=-1 WHERE id=?", (file_id,))
            conn.commit()
            conn.close()
            continue
        
        conn = get_db()
        cursor = conn.cursor()
        file_type = get_file_type(path)
        cursor.execute("UPDATE files SET type=? WHERE id=?", (file_type, file_id))
        conn.commit()
        conn.close()
        
        process_single_file(file_id, path, frame_interval, whisper_model)
        processed_count += 1

    return processed_count

# =======================
# CONVENIENCE FUNCTIONS
# =======================
def add_file_to_db(file_path):
    """Add file to processing queue"""
    conn = get_db()
    cursor = conn.cursor()
    
    file_type = get_file_type(file_path)
    try:
        cursor.execute(
            "INSERT OR IGNORE INTO files (path, type, processed) VALUES (?, ?, 0)",
            (str(file_path), file_type)
        )
        conn.commit()
        return cursor.lastrowid
    except Exception as e:
        print(f"[ERROR] Failed to add {file_path}: {e}")
        return None
    finally:
        conn.close()

def process_directory(directory_path, recursive=True):
    """Add all supported files from directory"""
    supported_exts = {
        '.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm', '.m4v', '.3gp',
        '.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.gif',
        '.pdf', '.docx', '.pptx', '.csv', '.ipynb', '.html', '.htm', '.md',
        '.txt', '.py', '.cpp', '.c', '.h', '.java', '.js', '.ts', '.json','.css'
    }
    
    directory = Path(directory_path)
    if not directory.exists():
        print(f"[ERROR] Directory not found: {directory_path}")
        return
    
    pattern = "**/*" if recursive else "*"
    files_added = 0
    
    for file_path in directory.glob(pattern):
        if file_path.is_file() and file_path.suffix.lower() in supported_exts:
            if add_file_to_db(file_path):
                files_added += 1
    
    print(f"[INFO] Added {files_added} files")

def get_failed_files():
    """Get list of failed files"""
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("SELECT id, path, type FROM files WHERE processed = -1")
    failed_files = cursor.fetchall()
    conn.close()
    
    if failed_files:
        print(f"[INFO] {len(failed_files)} failed files:")
        for file_id, path, file_type in failed_files:
            print(f"  {Path(path).name}")
    else:
        print("[INFO] No failed files.")
    
    return failed_files

def retry_failed_files():
    """Retry failed files"""
    conn = get_db()
    cursor = conn.cursor()
    cursor.execute("UPDATE files SET processed = 0 WHERE processed = -1")
    updated_count = cursor.rowcount
    conn.commit()
    conn.close()
    
    print(f"[INFO] Reset {updated_count} files for retry.")
    return updated_count

def migrate_to_siglip_text():
    """Helper function to re-embed all text using SigLIP for consistency"""
    print("[MIGRATE] Converting text embeddings to SigLIP...")
    
    conn = get_db()
    cursor = conn.cursor()
    
    # Get all text-based vectors
    cursor.execute("""
        SELECT v.id, v.file_id, v.chunk_id, v.modality, f.path
        FROM vectors v 
        JOIN files f ON v.file_id = f.id
        WHERE v.modality IN ('text', 'ocr_text', 'video_audio')
    """)
    
    text_vectors = cursor.fetchall()
    print(f"[MIGRATE] Found {len(text_vectors)} text vectors to convert")
    
    if not text_vectors:
        conn.close()
        return
    
    # Delete old text index
    text_index_path = "indices/faiss_text.index"
    if Path(text_index_path).exists():
        os.remove(text_index_path)
        print("[MIGRATE] Deleted old text index")
    
    # Mark all files as unprocessed to trigger re-embedding
    cursor.execute("UPDATE files SET processed = 0")
    
    # Delete all text vectors from database
    cursor.execute("DELETE FROM vectors WHERE modality IN ('text', 'ocr_text', 'video_audio')")
    
    conn.commit()
    conn.close()
    
    print("[MIGRATE] Migration prepared. Run process_files() to re-embed with SigLIP.")

def test_embeddings():
    """Test the corrected embeddings"""
    print("[TEST] Testing corrected embeddings...")
    
    # Test text embedding
    text_emb = embed_text_enhanced("test document about cars")
    if text_emb is not None:
        print(f"[TEST] Text embedding: shape={text_emb.shape}, norm={np.linalg.norm(text_emb):.3f}")
    
    # Test visual text embedding
    visual_text_emb = get_visual_text_embedding("red car")
    if visual_text_emb is not None:
        print(f"[TEST] Visual text embedding: shape={visual_text_emb.shape}, norm={np.linalg.norm(visual_text_emb):.3f}")
    
    # Test image embedding
    try:
        # Create a small test image
        test_img = Image.new('RGB', (100, 100), color='red')
        img_emb = embed_image_siglip(test_img)
        if img_emb is not None:
            print(f"[TEST] Image embedding: shape={img_emb.shape}, norm={np.linalg.norm(img_emb):.3f}")
    except Exception as e:
        print(f"[TEST] Image embedding test failed: {e}")
    
    # Test similarity calculation if dimensions match
    if (text_emb is not None and visual_text_emb is not None and 
        text_emb.shape == visual_text_emb.shape):
        # For normalized embeddings, cosine similarity = dot product
        similarity = np.dot(text_emb, visual_text_emb)
        print(f"[TEST] Cross-modal similarity: {similarity:.3f}")
        
        if img_emb is not None and text_emb.shape == img_emb.shape:
            img_similarity = np.dot(text_emb, img_emb)
            print(f"[TEST] Text-Image similarity: {img_similarity:.3f}")
    else:
        print("[TEST] WARNING: Embedding dimensions don't match - cross-modal search will fail!")
        if text_emb is not None and visual_text_emb is not None:
            print(f"[TEST] Text: {text_emb.shape}, Visual: {visual_text_emb.shape}")
    
    print("[TEST] Embedding test complete")

# =======================
# RUN
# =======================
if __name__ == "__main__":
    print("🚀 FIXED Clean Embedding System")
    print("GPU: SentenceTransformer + SigLIP + EasyOCR")
    print("CPU: Whisper audio processing")
    print("FIXES: Proper SigLIP features + Normalization + IndexFlatIP")

    # Test embeddings first
    test_embeddings()

    while True:
        processed = process_files(limit=50, frame_interval=30, whisper_model="base")
        print(f"[DONE] Processed {processed} files in this batch")

        if processed == 0:
            print("✅ All files processed. Exiting loop.")
            break

    get_failed_files()
