"""
output.py
---------
Handles semantic search and summarization.
- Loads FAISS indices (text + image)
- Embeds queries with SentenceTransformer / SigLIP / Ollama fallback
- Enhanced search combining semantic + filename ranking
- Summarization with optimized prompt generation
"""

import os
import re
import sqlite3
import subprocess
import platform
import mimetypes
import base64
import time
from pathlib import Path
from collections import deque, defaultdict
import fitz  # PyMuPDF
import easyocr
from PIL import Image
import io
import cv2
import tempfile

import numpy as np
import faiss
import torch
from sentence_transformers import SentenceTransformer
from transformers import AutoProcessor, AutoModel
from ollama import Client as OllamaClient

# Optional extractors
try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation  
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# Create config if not exists
try:
    import config
except ImportError:
    config_content = '''import os
from pathlib import Path

SQLITE_PATH = os.getenv("SQLITE_PATH", "files.db")
INDICES_DIR = Path("indices")
INDICES_DIR.mkdir(exist_ok=True)

# Search Configuration
DEFAULT_TOP_K = int(os.getenv("DEFAULT_TOP_K", "5"))
MAX_TOP_K = int(os.getenv("MAX_TOP_K", "10"))
TEXT_TRUNCATE = int(os.getenv("TEXT_TRUNCATE", "8000"))
MEMORY_TURNS = int(os.getenv("MEMORY_TURNS", "6"))

# Model Configuration  
TEXT_MODEL = os.getenv("TEXT_MODEL", "mistral:7b")
MULTI_MODEL = os.getenv("MULTI_MODEL", "llava")
PDF_PAGE_LIMIT = int(os.getenv("PDF_PAGE_LIMIT", "5"))
'''
    with open('config.py', 'w') as f:
        f.write(config_content)
    import config

# Configuration from config file and environment
DEFAULT_TOP_K = getattr(config, 'DEFAULT_TOP_K', 5)
MAX_TOP_K = getattr(config, 'MAX_TOP_K', 5)
TEXT_TRUNCATE = getattr(config, 'TEXT_TRUNCATE', 8000)
MEMORY_TURNS = getattr(config, 'MEMORY_TURNS', 6)
TEXT_MODEL = getattr(config, 'TEXT_MODEL', "mistral:7b")
MULTI_MODEL = getattr(config, 'MULTI_MODEL', "llava")
PDF_PAGE_LIMIT = getattr(config, 'PDF_PAGE_LIMIT', 5)

# Model initialization
device = "cuda" if torch.cuda.is_available() else "cpu"
print(f"Using device: {device}")

print("Initializing EasyOCR...")
easyocr_reader = easyocr.Reader(['en'], gpu=torch.cuda.is_available(), verbose=False)

print("Loading SigLIP model...")
try:
    siglip_model = AutoModel.from_pretrained("google/siglip-base-patch16-224").to(device)
    siglip_processor = AutoProcessor.from_pretrained("google/siglip-base-patch16-224")
    print("SigLIP loaded successfully")
except Exception as e:
    print(f"SigLIP failed to load: {e}")
    siglip_model = siglip_processor = None

try:
    text_embedder = SentenceTransformer("all-MiniLM-L6-v2")
    print("Text embedder loaded successfully")
except Exception as e:
    print(f"Text embedder failed: {e}")
    text_embedder = None

print("Connecting to Ollama...")
ollama = OllamaClient()

try:
    test_response = ollama.generate(model=TEXT_MODEL, prompt="test", options={"num_predict": 1})
    print("Ollama connected successfully")
except Exception as e:
    print(f"Ollama connection issue: {e}")

print("All models initialized!\n")

# Database helper
def get_db_conn():
    return sqlite3.connect(config.SQLITE_PATH)

# Embedding helpers
def embed_text(text):
    """Embed text using SentenceTransformer or Ollama fallback"""
    if not text or not text.strip():
        return None
    try:
        if text_embedder:
            e = text_embedder.encode(text, convert_to_numpy=True)
            return e.astype("float32")
    except Exception:
        pass
    try:
        res = ollama.embeddings(model=TEXT_MODEL, prompt=text)
        if isinstance(res, dict) and "embedding" in res:
            return np.array(res["embedding"], dtype="float32")
    except Exception:
        pass
    return None


def embed_visual(text):
    """Embed text for visual search using SigLIP (not CLIP)."""
    if siglip_model is None or siglip_processor is None:
        return None
    try:
        inputs = siglip_processor(text=[text], return_tensors="pt", padding=True).to(device)
        with torch.no_grad():
            features = siglip_model.get_text_features(**inputs)
        return features.cpu().numpy().flatten().astype("float32")
    except Exception as e:
        print(f"SigLIP embedding failed: {e}")
        return None

# Enhanced Search System
class ImprovedFileSearcher:
    def __init__(self, top_k=None):
        self.top_k = top_k or DEFAULT_TOP_K
        self.max_search_k = min(self.top_k * 4, 50)  # Dynamic search expansion
        self.text_index = None
        self.image_index = None
        self.load_indices()
    
    def load_indices(self):
        """Load FAISS indices if they exist"""
        text_idx_path = Path("indices/faiss_text.index")
        img_idx_path = Path("indices/faiss_image.index")
        
        if text_idx_path.exists():
            try:
                self.text_index = faiss.read_index(str(text_idx_path))
                print(f"Loaded text index with {self.text_index.ntotal} vectors")
            except Exception as e:
                print(f"Failed to load text index: {e}")
                
        if img_idx_path.exists():
            try:
                self.image_index = faiss.read_index(str(img_idx_path))
                print(f"Loaded image index with {self.image_index.ntotal} vectors")
            except Exception as e:
                print(f"Failed to load image index: {e}")
    
    def calculate_filename_score(self, query, filename):
        """Calculate relevance score based on filename matching"""
        query_words = set(re.findall(r'\w+', query.lower()))
        filename_words = set(re.findall(r'\w+', Path(filename).stem.lower()))
        
        if not query_words or not filename_words:
            return 0.0
        
        # Exact matches get highest score
        exact_matches = len(query_words.intersection(filename_words))
        
        # Partial matches (substring matching)
        partial_matches = 0
        for q_word in query_words:
            for f_word in filename_words:
                if len(q_word) >= 3 and (q_word in f_word or f_word in q_word):
                    partial_matches += 0.5
        
        # Normalize by query length with boost for high matches
        total_score = (exact_matches * 1.5 + partial_matches) / len(query_words)
        return min(total_score, 1.5)  # Allow boost above 1.0 for exact matches
    
    def detect_query_intent(self, query):
        """Detect what type of content user is looking for"""
        query_lower = query.lower()
        
        intent = {
            'file_type': None,
            'prefer_visual': False,
            'prefer_text': False,
            'specific_format': None
        }
        
        # File type detection
        if any(x in query_lower for x in ['image', 'photo', 'picture', 'pic', 'screenshot']):
            intent['file_type'] = 'image'
            intent['prefer_visual'] = True
        elif any(x in query_lower for x in ['video', 'movie', 'clip', 'recording']):
            intent['file_type'] = 'video'
            intent['prefer_visual'] = True
        elif any(x in query_lower for x in ['document', 'pdf', 'doc', 'text', 'file']):
            intent['file_type'] = 'document'
            intent['prefer_text'] = True
        elif any(x in query_lower for x in ['.py', 'python', 'code', 'script']):
            intent['file_type'] = 'document'
            intent['specific_format'] = 'code'
            intent['prefer_text'] = True
        
        return intent
    
    def enhanced_search(self, query, top_k=None):
        """Enhanced search with multiple ranking factors"""
        if top_k:
            self.top_k = min(top_k, MAX_TOP_K)
        
        results = []
        query_lower = query.lower()
        intent = self.detect_query_intent(query)
        
        conn = get_db_conn()
        cursor = conn.cursor()
        
        # Get all files for filename fallback
        cursor.execute("SELECT id, path, type FROM files")
        all_files = {row[0]: (row[1], row[2]) for row in cursor.fetchall()}
        
        # SEMANTIC SEARCH - Text Index
        text_results = self._search_text_index(query, cursor, intent)
        
        # SEMANTIC SEARCH - Image Index  
        image_results = self._search_image_index(query, cursor, intent)
        
        # FILENAME FALLBACK SEARCH - Critical for exact matches
        filename_results = self._search_by_filename(query, all_files, intent)
        
        conn.close()
        
        # COMBINE AND RANK RESULTS
        all_results = text_results + image_results + filename_results
        unique_results = self._deduplicate_and_rank(all_results, intent)
        
        print(f"Search '{query}': {len(text_results)} text, {len(image_results)} visual, {len(filename_results)} filename matches")
        return unique_results[:self.top_k]
    
    def _search_text_index(self, query, cursor, intent):
        """Search text index with scoring"""
        results = []
        if self.text_index is None:
            return results
            
        text_emb = embed_text(query)
        if text_emb is None:
            return results
        
        try:
            search_k = min(self.max_search_k, self.text_index.ntotal)
            D, I = self.text_index.search(np.array([text_emb]), search_k)
            
            cursor.execute("""
                SELECT v.id, v.file_id, f.path, f.type, v.modality, v.chunk_id
                FROM vectors v JOIN files f ON v.file_id = f.id
                WHERE v.modality IN ('text', 'video_audio', 'ocr_text')
                ORDER BY v.id
            """)
            text_vectors = cursor.fetchall()
            
            for dist, idx in zip(D[0], I[0]):
                if idx != -1 and idx < len(text_vectors):
                    _, file_id, path, ftype, modality, chunk_id = text_vectors[idx]
                    
                    # Apply file type filter based on intent
                    if intent['file_type'] and ftype != intent['file_type']:
                        continue
                    
                    semantic_score = float(1.0 / (1.0 + float(dist)))
                    filename_score = self.calculate_filename_score(query, path)
                    
                    # Modality-specific scoring
                    if modality == 'text':
                        weight_semantic, weight_filename = 0.75, 0.25
                    elif modality == 'ocr_text':
                        weight_semantic, weight_filename = 0.85, 0.15  # OCR text is highly relevant
                    else:  # video_audio
                        weight_semantic, weight_filename = 0.65, 0.35
                    
                    # Intent-based scoring adjustments
                    if intent['prefer_text'] and modality in ['text', 'ocr_text']:
                        weight_semantic += 0.1
                    
                    final_score = semantic_score * weight_semantic + filename_score * weight_filename
                    
                    results.append({
                        'file_id': file_id,
                        'path': path,
                        'type': ftype,
                        'modality': modality,
                        'score': final_score,
                        'semantic_score': semantic_score,
                        'filename_score': filename_score,
                        'search_type': 'semantic_text'
                    })
                    
        except Exception as e:
            print(f"Text search error: {e}")
            
        return results
    
    def _search_image_index(self, query, cursor, intent):
        """Search image index with visual prioritization"""
        results = []
        if self.image_index is None:
            return results
            
        visual_emb = embed_visual(query)
        if visual_emb is None:
            return results
        
        try:
            search_k = min(self.max_search_k, self.image_index.ntotal)
            D, I = self.image_index.search(np.array([visual_emb]), search_k)
            
            cursor.execute("""
                SELECT v.id, v.file_id, f.path, f.type, v.modality, v.chunk_id
                FROM vectors v JOIN files f ON v.file_id = f.id
                WHERE v.modality IN ('image', 'video_frame')
                ORDER BY v.id
            """)
            visual_vectors = cursor.fetchall()
            
            for dist, idx in zip(D[0], I[0]):
                if idx != -1 and idx < len(visual_vectors):
                    _, file_id, path, ftype, modality, chunk_id = visual_vectors[idx]
                    
                    if intent['file_type'] and ftype != intent['file_type']:
                        continue
                    
                    semantic_score = float(1.0 / (1.0 + float(dist)))
                    filename_score = self.calculate_filename_score(query, path)
                    
                    # CRITICAL: Prioritize images over video frames for visual content
                    if modality == 'image':
                        weight_semantic, weight_filename = 0.80, 0.20
                        priority_boost = 0.15 if intent['prefer_visual'] else 0.05
                    else:  # video_frame
                        weight_semantic, weight_filename = 0.70, 0.30
                        priority_boost = 0.0
                    
                    final_score = semantic_score * weight_semantic + filename_score * weight_filename + priority_boost
                    
                    results.append({
                        'file_id': file_id,
                        'path': path,
                        'type': ftype,
                        'modality': modality,
                        'score': final_score,
                        'semantic_score': semantic_score,
                        'filename_score': filename_score,
                        'search_type': 'semantic_visual'
                    })
                    
        except Exception as e:
            print(f"Image search error: {e}")
            
        return results
    
    def _search_by_filename(self, query, all_files, intent):
        """Filename-based search for exact matches (critical for cases like 'badminton.mp4')"""
        results = []
        query_words = set(re.findall(r'\w+', query.lower()))
        
        for file_id, (path, ftype) in all_files.items():
            if intent['file_type'] and ftype != intent['file_type']:
                continue
            
            filename_score = self.calculate_filename_score(query, path)
            
            # More generous threshold for filename matches
            if filename_score >= 0.2:
                # Boost exact filename matches significantly
                if filename_score >= 1.0:
                    boosted_score = filename_score * 0.9  # High score for exact matches
                else:
                    boosted_score = filename_score * 0.5  # Lower for partial matches
                
                results.append({
                    'file_id': file_id,
                    'path': path,
                    'type': ftype,
                    'modality': 'filename_match',
                    'score': boosted_score,
                    'semantic_score': 0.0,
                    'filename_score': filename_score,
                    'search_type': 'filename'
                })
        
        return results
    
    def _deduplicate_and_rank(self, all_results, intent):
        """Deduplicate and rank results with intent-based scoring"""
        # Group by file to avoid duplicates
        file_groups = defaultdict(list)
        for result in all_results:
            file_groups[result['path']].append(result)
        
        unique_results = []
        for path, file_results in file_groups.items():
            if len(file_results) == 1:
                unique_results.append(file_results[0])
            else:
                # Multiple modalities for same file - combine intelligently
                best_result = max(file_results, key=lambda x: x['score'])
                
                # Enhance score based on modality diversity
                modalities = [r['modality'] for r in file_results]
                total_semantic = sum(r['semantic_score'] for r in file_results) / len(file_results)
                best_filename = max(r['filename_score'] for r in file_results)
                
                # Multi-modal bonus
                if len(set(modalities)) > 1:
                    multimodal_bonus = 0.1
                    best_result['score'] = total_semantic * 0.6 + best_filename * 0.3 + multimodal_bonus
                    best_result['search_type'] = 'multi_modal'
                
                unique_results.append(best_result)
        
        # Final ranking with intent consideration
        for result in unique_results:
            # Boost results that match user intent
            if intent['prefer_visual'] and result['modality'] in ['image', 'video_frame']:
                result['score'] *= 1.1
            elif intent['prefer_text'] and result['modality'] in ['text', 'ocr_text']:
                result['score'] *= 1.1
            elif result['search_type'] == 'filename' and result['filename_score'] > 0.8:
                result['score'] *= 1.2  # Boost high filename matches
        
        # Sort by final score
        unique_results.sort(key=lambda x: x['score'], reverse=True)
        return unique_results

# Enhanced Summarization System
class ImprovedFileSummarizer:
    def __init__(self, text_model=None, multi_model=None):
        self.text_model = text_model or TEXT_MODEL
        self.multi_model = multi_model or MULTI_MODEL
    
    def detect_query_type(self, query):
        """Detect what kind of analysis user wants"""
        query_lower = query.lower()
        
        if any(word in query_lower for word in ['extract', 'find', 'get', 'list', 'show', 'what']):
            return 'extraction'
        elif any(word in query_lower for word in ['summarize', 'summary', 'overview', 'explain']):
            return 'summary'
        elif any(word in query_lower for word in ['analyze', 'analysis', 'details']):
            return 'analysis'
        else:
            return 'general'
    
    def create_enhanced_prompt(self, query, content, file_info, query_type):
        """Create optimized prompts based on query type"""
        file_name = Path(file_info.get('name', 'document')).stem
        
        base_context = f"""Document: {file_name}
File Size: {file_info.get('size', 'unknown')}
Content Type: {file_info.get('mime', 'unknown')}

DOCUMENT CONTENT:
---
{content}
---

User Query: "{query}"
"""

        if query_type == 'extraction':
            prompt = f"""{base_context}

TASK: Extract ALL relevant information that matches the user's query. Be comprehensive and include:

1. ALL specific details, numbers, names, dates, addresses, IDs mentioned
2. Any structured information (forms, tables, lists)
3. Related information that might be useful
4. Maintain exact formatting and context

IMPORTANT: Do not summarize - provide complete, verbatim information as it appears in the document. Include everything relevant, even if it seems redundant."""

        elif query_type == 'summary':
            prompt = f"""{base_context}

TASK: Provide a comprehensive summary that includes:

1. Main topics and key information points
2. Important details, numbers, dates, and facts
3. Any structured data (IDs, names, addresses, forms)
4. Context and relationships between information
5. Document purpose and key takeaways

Organize the information clearly and don't omit important details."""

        elif query_type == 'analysis':
            prompt = f"""{base_context}

TASK: Provide detailed analysis including:

1. Thorough examination of all content
2. Identification of key information and patterns  
3. Explanation of document structure and purpose
4. Detailed breakdown of important elements
5. Context and significance of the information

Be analytical but stick to facts present in the document."""

        else:  # general
            prompt = f"""{base_context}

TASK: Respond to the user's query comprehensively:

1. Address their specific question thoroughly
2. Include all relevant details from the document
3. Provide context and explanation as needed
4. Don't make assumptions beyond what's in the document
5. If the document contains forms or structured data, present it clearly

Ensure your response is complete and helpful."""

        return prompt
    
    def summarize_with_context(self, query, file_path, context=""):
        """Enhanced file summarization with better accuracy and restricted Llava usage"""
        try:
            info = get_file_info(file_path)
            if not info:
                return "Could not access file information."
            
            query_type = self.detect_query_type(query)
            
            # CRITICAL: Check if this is a standalone image or video file
            is_standalone_visual = info.get('is_image') or info.get('is_video')
            
            # For TEXT-BASED FILES (PDFs, DOCX, TXT, etc.) - Use TEXT MODEL ONLY
            if not is_standalone_visual:
                print(f"Processing text-based file with text model: {info['name']}")
                
                # Try direct text extraction first for better accuracy
                direct_content = None
                if info.get('is_pdf') or info.get('is_docx') or info.get('is_text'):
                    try:
                        direct_content = read_file_local(file_path)
                    except:
                        pass
                
                # Use direct extraction if available and substantial
                if direct_content and len(direct_content.strip()) > 50:
                    # Truncate if too long but preserve key information
                    if len(direct_content) > TEXT_TRUNCATE:
                        # Try to find natural breaking points
                        sentences = re.split(r'[.!?]+', direct_content)
                        truncated = ""
                        for sentence in sentences:
                            if len(truncated + sentence) > TEXT_TRUNCATE:
                                break
                            truncated += sentence + ". "
                        direct_content = truncated.strip()
                    
                    enhanced_prompt = self.create_enhanced_prompt(
                        query, direct_content, info, query_type
                    )
                    
                    messages = [
                        {"role": "system", "content": "You are a thorough document analysis assistant. Provide complete, accurate information based exactly on the document content. Do not omit important details."},
                        {"role": "user", "content": enhanced_prompt}
                    ]
                    
                    result = chat_with_ollama_messages(messages, model=self.text_model)
                    if result and not result.startswith("[Model error"):
                        return f"Analysis of {info['name']}:\n\n{result}"
                
                # Fallback for PDFs - use OCR + text model (NOT Llava)
                if info.get('is_pdf'):
                    print("  Falling back to OCR + text model for PDF")
                    return self._process_pdf_with_ocr_only(query, file_path, info, query_type)
                
                # For other text files that failed direct extraction
                return f"Could not extract text content from {info['name']}"
            
            # FOR STANDALONE IMAGE/VIDEO FILES ONLY - Use Llava
            else:
                print(f"Processing standalone visual file with Llava: {info['name']}")
                result = send_file_to_ollama_multimodal(
                    query, file_path, context, self.multi_model, self.text_model
                )
                return result if result else f"Could not analyze {info['name']}"
            
        except Exception as e:
            return f"Error processing file: {e}"
    
    def _process_pdf_with_ocr_only(self, query, file_path, info, query_type):
        """Process PDF using OCR + text model only (no Llava)"""
        try:
            print("  Extracting text from PDF using OCR...")
            ocr_text = ""
            
            doc = fitz.open(str(file_path))
            for page_num, page in enumerate(doc):
                # Get text first (if available)
                page_text = page.get_text()
                if page_text.strip():
                    ocr_text += f"Page {page_num + 1} (Direct Text):\n{page_text}\n\n"
                else:
                    # Only do OCR if no direct text
                    print(f"    OCR processing page {page_num + 1}")
                    pix = page.get_pixmap(dpi=300)
                    mode = "RGB" if pix.n < 4 else "RGBA"
                    img = Image.frombytes(mode, (pix.width, pix.height), pix.samples)
                    ocr_result = run_easyocr_on_image(img)
                    if ocr_result:
                        ocr_text += f"Page {page_num + 1} (OCR):\n{ocr_result}\n\n"
                
                if page_num + 1 >= PDF_PAGE_LIMIT:
                    break
            
            doc.close()
            
            if not ocr_text.strip():
                return f"No text could be extracted from {info['name']}"
            
            # Create enhanced prompt for OCR content
            enhanced_prompt = self.create_enhanced_prompt(
                query, ocr_text, info, query_type
            )
            
            messages = [
                {"role": "system", "content": "You are analyzing a document based on extracted text content. Provide complete, accurate information based on the text provided."},
                {"role": "user", "content": enhanced_prompt}
            ]
            
            result = chat_with_ollama_messages(messages, model=self.text_model)
            return f"Analysis of {info['name']}:\n\n{result}" if result else f"Could not analyze extracted text from {info['name']}"
            
        except Exception as e:
            return f"Error processing PDF with OCR: {e}"



# File utilities (keeping existing functions)
def open_file(path):
    """Open file with system default application"""
    p = Path(path)
    if not p.exists():
        return f"File not found: {path}"
    try:
        system = platform.system()
        if system == "Windows":
            os.startfile(str(p))
        elif system == "Darwin":
            subprocess.run(["open", str(p)])
        else:
            subprocess.run(["xdg-open", str(p)])
        return f"Opened: {p.name}"
    except Exception as e:
        return f"Error opening file: {e}"

def get_file_info(path):
    """Get file information"""
    p = Path(path)
    if not p.exists():
        return None
    size = p.stat().st_size
    size_str = f"{size/1024:.1f} KB" if size < 1e6 else f"{size/1e6:.1f} MB"
    ext = p.suffix.lower()
    mime, _ = mimetypes.guess_type(str(p))
    return {
        "name": p.name,
        "size": size_str,
        "ext": ext,
        "mime": mime or "unknown",
        "is_text": ext in ['.txt', '.py', '.js', '.md', '.json', '.html', '.csv'],
        "is_pdf": ext == '.pdf',
        "is_docx": ext == '.docx',
        "is_pptx": ext == '.pptx',
        "is_image": ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff'],
        "is_video": ext in ['.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv']
    }

def read_file_local(path):
    """Read file content locally"""
    try:
        p = str(path)
        if p.lower().endswith(".pdf"):
            text = ""
            with fitz.open(p) as doc:
                for page in doc[:PDF_PAGE_LIMIT]:  # Use configurable limit
                    text += page.get_text()
            return text.strip() or "[No text found]"
        if p.lower().endswith(".docx") and DOCX_AVAILABLE:
            return "\n".join([para.text for para in Document(p).paragraphs]).strip() or "[No text found]"
        if p.lower().endswith(".pptx") and PPTX_AVAILABLE:
            texts = []
            for slide in Presentation(p).slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts).strip() or "[No text found]"
        with open(p, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception as e:
        return f"[Error extracting file: {e}]"
    
def extract_video_frames_for_analysis(video_path, max_frames=5, frame_interval=30):
    """
    Extract frames from video for Llava analysis using the same sampling strategy as embedding.
    Uses every nth frame to avoid overwhelming the model while maintaining temporal coverage.
    
    Args:
        video_path: Path to video file
        max_frames: Maximum number of frames to extract for analysis
        frame_interval: Extract every nth frame (should match embedding strategy)
    
    Returns:
        List of base64 encoded frames
    """
    import cv2
    import tempfile
    
    frames_b64 = []
    
    try:
        cap = cv2.VideoCapture(str(video_path))
        if not cap.isOpened():
            print(f"Could not open video: {video_path}")
            return frames_b64
        
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        fps = cap.get(cv2.CAP_PROP_FPS)
        duration = total_frames / fps if fps > 0 else 0
        
        print(f"  Video info: {total_frames} frames, {fps:.1f} fps, {duration:.1f}s duration")
        
        # Calculate which frames to extract
        # Start from frame_interval to match embedding sampling
        frame_numbers = []
        current_frame = frame_interval  # Start at 30, then 60, 90, etc.
        
        while current_frame < total_frames and len(frame_numbers) < max_frames:
            frame_numbers.append(current_frame)
            current_frame += frame_interval  # Jump by interval (30 frames)
        
        print(f"  Extracting frames: {frame_numbers}")
        
        for frame_num in frame_numbers:
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_num)
            ret, frame = cap.read()
            
            if ret:
                # Convert BGR to RGB
                frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                
                # Convert to PIL Image
                pil_image = Image.fromarray(frame_rgb)
                
                # Resize if too large (optional optimization)
                max_size = 1024
                if max(pil_image.size) > max_size:
                    pil_image.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
                
                # Convert to base64 using in-memory buffer (avoids temp file permissions)
                img_buffer = io.BytesIO()
                pil_image.save(img_buffer, format='JPEG', quality=85)
                img_buffer.seek(0)
                frame_b64 = base64.b64encode(img_buffer.getvalue()).decode('utf-8')
                frames_b64.append(frame_b64)
                img_buffer.close()
                
                timestamp = frame_num / fps if fps > 0 else frame_num
                print(f"    Frame {frame_num} (t={timestamp:.1f}s) extracted")
            else:
                print(f"    Could not read frame {frame_num}")
        
        cap.release()
        print(f"  Successfully extracted {len(frames_b64)} frames")
        
    except Exception as e:
        print(f"Error extracting video frames: {e}")
        if 'cap' in locals():
            cap.release()
    
    return frames_b64

# Ollama response handling (keeping existing functions)
def extract_ollama_text(response):
    """Extract clean text from Ollama response"""
    if response is None:
        return ""
    if isinstance(response, dict):
        if "message" in response:
            msg = response["message"]
            if isinstance(msg, dict) and "content" in msg:
                return str(msg["content"]).strip()
            if hasattr(msg, "content"):
                return str(getattr(msg, "content", "")).strip()
        if "response" in response and isinstance(response["response"], str):
            return response["response"].strip()
        if "text" in response and isinstance(response["text"], str):
            return response["text"].strip()
        if "content" in response and isinstance(response["content"], str):
            return response["content"].strip()
    if hasattr(response, "message") and hasattr(response.message, "content"):
        return str(response.message.content).strip()
    if hasattr(response, "content"):
        return str(response.content).strip()
    s = str(response)
    m = re.search(r"content=['\"]([^'\"]+)['\"]", s)
    if m:
        return m.group(1).strip()
    m = re.search(r"response=['\"]([^'\"]+)['\"]", s)
    if m:
        return m.group(1).strip()
    return s.strip()

def clean_model_text(text):
    """Clean model response text"""
    if not text:
        return ""
    t = str(text).strip()
    t = re.sub(r'(?mi)^\s*model=.*$', '', t)
    t = re.sub(r'(?mi)^\s*created_at=.*$', '', t)
    t = re.sub(r'(?mi)^\s*done=.*$', '', t)
    t = re.sub(r'(?mi)^\s*done_reason=.*$', '', t)
    t = re.sub(r'(?mi)^\s*total_duration=.*$', '', t)
    t = re.sub(r'(?mi)^\s*load_duration=.*$', '', t)
    t = re.sub(r'(?mi)^\s*prompt_eval_count=.*$', '', t)
    t = re.sub(r'(?mi)^\s*prompt_eval_duration=.*$', '', t)
    t = re.sub(r'(?mi)^\s*eval_count=.*$', '', t)
    t = re.sub(r'(?mi)^\s*eval_duration=.*$', '', t)
    t = re.sub(r'(?mi)^\s*message=.*$', '', t)
    t = re.sub(r'(?mi)^\s*thinking=.*$', '', t)
    t = re.sub(r'(?mi)^\s*images=.*$', '', t)
    t = re.sub(r'(?mi)^\s*tool_name=.*$', '', t)
    t = re.sub(r'(?mi)^\s*tool_calls=.*$', '', t)
    t = re.sub(r'\n\s*\n\s*', '\n\n', t)
    return t.strip()

def chat_with_ollama_messages(messages, model=None):
    """Chat with Ollama using messages format"""
    model = model or TEXT_MODEL
    try:
        res = ollama.chat(model=model, messages=messages)
        out = extract_ollama_text(res)
        return clean_model_text(out)
    except Exception:
        try:
            prompt = ""
            for m in messages:
                role = m.get("role")
                content = m.get("content", "")
                prompt += f"{role.upper()}: {content}\n"
            res2 = ollama.generate(model=model, prompt=prompt)
            out2 = extract_ollama_text(res2)
            return clean_model_text(out2)
        except Exception as e:
            return f"[Model error: {e}]"

# OCR and multimodal processing
def run_easyocr_on_image(img, ocr_langs=['en']):
    """Run EasyOCR on a PIL Image object"""
    try:
        img_np = np.array(img)
        
        if ocr_langs != ['en']:
            reader = easyocr.Reader(ocr_langs, gpu=torch.cuda.is_available(), verbose=False)
        else:
            reader = easyocr_reader
        
        results = reader.readtext(img_np)
        extracted_texts = [result[1] for result in results if result[2] > 0.1]
        
        return ' '.join(extracted_texts).strip()
        
    except Exception as e:
        print(f"EasyOCR error: {e}")
        return ""

def encode_image_to_base64(image_path):
    """Encode image to base64 for multimodal models"""
    try:
        with open(image_path, "rb") as image_file:
            return base64.b64encode(image_file.read()).decode('utf-8')
    except Exception as e:
        print(f"Error encoding image: {e}")
        return None

def send_file_to_ollama_multimodal(query, file_path, context="", multi_model=None, text_model=None):
    """Enhanced multimodal processing with proper video frame extraction"""
    multi_model = multi_model or MULTI_MODEL
    text_model = text_model or TEXT_MODEL
    
    info = get_file_info(file_path)
    if not info:
        return "File not found."

    # CRITICAL CHECK: Only allow standalone image/video files
    if not (info.get("is_image") or info.get("is_video")):
        return f"Error: Multimodal processing is restricted to standalone image/video files only. {info['name']} is not supported for Llava processing."

    print(f"Processing standalone visual file with Llava: {info['name']} ({info['size']})")

    images_b64 = []
    ocr_outputs = []
    frame_info = ""

    # Process standalone image files
    if info.get("is_image"):
        print("  Processing standalone image...")
        b64 = encode_image_to_base64(file_path)
        if b64:
            images_b64.append(b64)
        try:
            img = Image.open(file_path)
            text = run_easyocr_on_image(img)
            if text:
                ocr_outputs.append(text)
            print(f"  OCR extracted: {len(text) if text else 0} characters")
        except Exception as e:
            print(f"Image OCR failed: {e}")

    # Process standalone video files with frame sampling
    elif info.get("is_video"):
        print("  Processing standalone video with frame sampling...")
        
        # Extract frames using the same sampling strategy as embedding
        video_frames = extract_video_frames_for_analysis(
            file_path, 
            max_frames=5,  # Limit frames to avoid overwhelming Llava
            frame_interval=60  # Match embedding frame interval
        )
        
        if video_frames:
            images_b64.extend(video_frames)
            frame_info = f"Extracted {len(video_frames)} frames sampled at 30-frame intervals for analysis."
            print(f"  {frame_info}")
        else:
            print("  No frames could be extracted from video")
            return f"Could not extract frames from video {info['name']}"

    # Join OCR text
    ocr_text_joined = "\n\n".join(ocr_outputs).strip()
    if ocr_text_joined:
        print(f"  Total OCR text length: {len(ocr_text_joined)}")

    # Create enhanced prompt for video analysis
    if info.get("is_video"):
        visual_prompt = f"""
You are analyzing a video file. The user asked: "{query}"

VIDEO INFORMATION:
- File: {info['name']} ({info['size']})
- {frame_info}
- Frames are sampled at regular intervals to represent the video content

OCR EXTRACTED TEXT (from video frames, if any):
---
{ocr_text_joined if ocr_text_joined else "[No OCR text found in video frames]"}
---

INSTRUCTIONS:
1. Analyze the visual content across all provided video frames
2. Look for patterns, changes, or consistent elements across frames
3. Identify objects, people, actions, text, scenes, and settings
4. Note any temporal progression or movement if visible
5. Use OCR text to supplement your visual understanding
6. Answer the user's question comprehensively based on video content
7. If this appears to be a sports video (like badminton), focus on:
   - Players and their actions
   - Game elements (court, equipment, etc.)
   - Scores or game information visible
   - Playing techniques or movements

Please provide a comprehensive analysis of this video content based on the sampled frames.
"""
    else:
        # Image analysis prompt (unchanged)
        visual_prompt = f"""
You are analyzing a standalone image file. The user asked: "{query}"

OCR EXTRACTED TEXT (if any):
---
{ocr_text_joined if ocr_text_joined else "[No OCR text found in image]"}
---

INSTRUCTIONS:
1. Analyze the visual content in the image
2. Use any OCR text to supplement your visual understanding
3. Answer the user's question comprehensively
4. Focus on visual elements: objects, people, text, scenes, colors, composition
5. If there's text in the image, include it in your analysis
6. Be thorough but concise

Please provide a comprehensive analysis of this visual content.
"""

    # Send to multimodal model (Llava)
    print("  Analyzing with Llava multimodal model...")
    try:
        if images_b64:
            # Use Llava for visual analysis
            response = ollama.generate(
                model=multi_model,
                prompt=visual_prompt,
                images=images_b64
            )
            result = extract_ollama_text(response)
            print("  Visual analysis complete")
            
            # Add frame info to result for videos
            if info.get("is_video") and frame_info:
                result = f"[{frame_info}]\n\n{result}"
            
            return clean_model_text(result) if result else f"Could not analyze {info['name']}"
        else:
            return f"No visual content could be processed from {info['name']}"
            
    except Exception as e:
        print(f"Llava analysis failed: {e}")
        if ocr_text_joined:
            return f"Llava failed, but extracted OCR text:\n{ocr_text_joined}"
        else:
            return f"Could not process visual file {info['name']} - {e}"

# Action detection and file resolution
def decide_action(query):
    """Determine what action to take based on query"""
    q = query.lower().strip()
    if any(k in q for k in ["file", "read", "read-ocr", "read-raw", "summarize", "analyze", "explain", "extract"]) and re.search(r'\d+', q):
        return "file_task"
    search_terms = ['find', 'search', 'show me', 'get me', 'list', 'locate', 'files']
    if any(term in q for term in search_terms) and not q.startswith(("how", "what", "why", "explain")):
        return "search"
    return "chat"

def resolve_file_reference(query, last_results):
    """Resolve file references in queries"""
    if not last_results:
        return None
        
    qlow = query.lower()
    m = re.search(r'\bfile\s+(\d+)\b', qlow)
    if m:
        idx = int(m.group(1)) - 1
        if 0 <= idx < len(last_results):
            return last_results[idx]["path"]
    m2 = re.search(r'\b(?:read|read-ocr|read-raw|summarize|analyze|explain|extract)\s+(\d+)\b', qlow)
    if m2:
        idx = int(m2.group(1)) - 1
        if 0 <= idx < len(last_results):
            return last_results[idx]["path"]
    if re.match(r'^\s*\d+\s*, query'):
        idx = int(query.strip()) - 1
        if 0 <= idx < len(last_results):
            return last_results[idx]["path"]
    return None

def format_search_results(results, query):
    """Format search results with better information"""
    if not results:
        return f"No files found for '{query}'. Try different keywords or check if files are indexed."
    
    lines = [f"Found {len(results)} most relevant files:"]
    
    for i, result in enumerate(results, 1):
        name = Path(result["path"]).name
        ftype = result.get("type", "unknown")
        score = result.get("score", 0.0)
        modality = result.get("modality", "unknown")
        search_type = result.get("search_type", "unknown")
        
        # Enhanced file type emojis
        emoji_map = {
            'image': '🖼️', 'video': '🎥', 'document': '📄',
            'py': '🐍', 'js': '🌐', 'pdf': '📕', 'docx': '📘'
        }
        emoji = emoji_map.get(ftype, '📄')
        
        # Show search method for debugging
        search_indicator = ""
        if search_type == "filename":
            search_indicator = " [filename match]"
        elif search_type == "multi_modal":
            search_indicator = " [multi-modal]"
        elif "semantic" in search_type:
            search_indicator = f" [{modality}]"
        
        lines.append(f"{i}. {emoji} {name} - Score: {score:.3f}{search_indicator}")
    
    lines.append(f"\nType number (1-{len(results)}) to open, or 'read N' to analyze.")
    lines.append("Note: Llava multimodal analysis is only available for standalone image/video files.")
    return "\n".join(lines)

# Enhanced Bot class
class ImprovedFileChatbot:
    def __init__(self, memory_turns=None, top_k=None):
        self.memory_turns = memory_turns or MEMORY_TURNS
        self.top_k = top_k or DEFAULT_TOP_K
        self.history = deque(maxlen=self.memory_turns * 2)
        self.last_results = []
        self.searcher = ImprovedFileSearcher(top_k=self.top_k)
        self.summarizer = ImprovedFileSummarizer()

    def set_top_k(self, k):
        """Dynamically set number of results to return"""
        self.top_k = min(k, MAX_TOP_K)
        self.searcher.top_k = self.top_k

    def add_history(self, user_msg, bot_msg):
        self.history.append(("user", user_msg))
        self.history.append(("assistant", bot_msg))

    def handle_search(self, query):
        """Enhanced search with improved ranking"""
        print(f"Searching with improved algorithm (top {self.top_k})...")
        results = self.searcher.enhanced_search(query, top_k=self.top_k)
        self.last_results = results
        
        response = format_search_results(results, query)
        self.add_history(query, response)
        return response

    def handle_file_task(self, query):
        """Enhanced file processing with restricted Llava usage"""
        file_path = resolve_file_reference(query, self.last_results)
        if not file_path:
            return "Please search for files first and then use 'read N' or provide a file path."

        if not Path(file_path).exists():
            return f"File not found: {file_path}"

        info = get_file_info(file_path)
        if not info:
            return "Could not access file information."

        print(f"Processing: {info['name']}")

        # Show processing method based on file type
        is_standalone_visual = info.get('is_image') or info.get('is_video')
        if is_standalone_visual:
            print(f"  → Will use Llava for standalone visual file")
        else:
            print(f"  → Will use text model for document/text file")

        context = "\n".join([f"{r[0]}: {r[1][:100]}" for r in list(self.history)[-6:]])
        is_read_ocr = bool(re.search(r'\bread-ocr\b', query.lower()))

        start_time = time.time()
        
        try:
            if is_read_ocr:
                # Return raw OCR only (for any file that supports it)
                if info.get('is_pdf') or info.get('is_image'):
                    if info.get('is_pdf'):
                        ocr_text = ""
                        doc = fitz.open(str(file_path))
                        for page_num, page in enumerate(doc):
                            pix = page.get_pixmap(dpi=300)
                            mode = "RGB" if pix.n < 4 else "RGBA"
                            img = Image.frombytes(mode, (pix.width, pix.height), pix.samples)
                            text = run_easyocr_on_image(img)
                            if text:
                                ocr_text += f"Page {page_num + 1}: {text}\n\n"
                            if page_num + 1 >= PDF_PAGE_LIMIT:
                                break
                        doc.close()
                    else:  # image
                        img = Image.open(file_path)
                        ocr_text = run_easyocr_on_image(img)
                    
                    result = f"Raw OCR output from {info['name']}:\n\n{ocr_text or '[No OCR text found]'}"
                else:
                    result = "OCR is only available for PDF and image files."
            else:
                # Full analysis using improved summarizer with restricted Llava
                result = self.summarizer.summarize_with_context(query, file_path, context)
                if not result:
                    result = f"Unable to process {info['name']}. File might be corrupted or unsupported."
                    
        except Exception as e:
            result = f"Error processing file: {e}"

        processing_time = time.time() - start_time
        print(f"  Completed in {processing_time:.1f}s")

        # Format response
        processing_method = "Llava multimodal" if is_standalone_visual else "Text model + OCR"
        bot_msg = f"{info['name']}\nProcessed: {info['size']} using {processing_method}, completed in {processing_time:.1f}s\n\n{result}"
        self.add_history(query, bot_msg)
        return bot_msg

    def process(self, query):
        """Main query processing"""
        q = query.strip()
        if not q:
            return ""

        # Handle direct file selection by number
        if q.isdigit() and self.last_results:
            idx = int(q) - 1
            if 0 <= idx < len(self.last_results):
                return open_file(self.last_results[idx]["path"])
            return "Invalid selection number."

        # Handle top-k configuration
        topk_match = re.search(r'\btop[-_\s]*(\d+)\b', q.lower())
        if topk_match:
            new_k = int(topk_match.group(1))
            self.set_top_k(new_k)
            # Remove the top-k instruction from query
            q = re.sub(r'\btop[-_\s]*\d+\b', '', q, flags=re.IGNORECASE).strip()
            if not q:
                return f"Set to return top {self.top_k} results. What would you like to search for?"

        action = decide_action(q)
        if action == "search":
            return self.handle_search(q)
        elif action == "file_task":
            return self.handle_file_task(q)
        else:
            # General chat
            context = "\n".join([f"{r[0]}: {r[1]}" for r in list(self.history)[-8:]])
            messages = [{"role": "system", "content": "You are a helpful file management assistant."}]
            if context:
                messages.append({"role": "assistant", "content": context})
            messages.append({"role": "user", "content": q})
            resp = chat_with_ollama_messages(messages, model=TEXT_MODEL)
            self.add_history(q, resp)
            return resp

# Directory scanning functions (integrates with embedder)
def scan_directory(directory_path):
    """Scan directory and add files to database using embedder"""
    try:
        from embedder import add_file_to_db
        
        directory = Path(directory_path)
        if not directory.exists():
            print(f"Directory not found: {directory_path}")
            return 0
        
        added_count = 0
        supported_extensions = {
            '.pdf', '.docx', '.pptx', '.txt', '.py', '.js', '.md', '.json', '.html', '.csv', '.cpp', '.c', '.h', '.java',
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.tif',
            '.mp4', '.avi', '.mov', '.mkv', '.wmv', '.flv', '.webm', '.m4v'
        }
        
        print(f"Scanning directory: {directory_path}")
        
        for file_path in directory.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in supported_extensions:
                if add_file_to_db(str(file_path)):
                    added_count += 1
                    if added_count % 50 == 0:
                        print(f"Added {added_count} files...")
        
        print(f"Scan complete. Added {added_count} files to database.")
        return added_count
        
    except ImportError:
        print("Error: embedder module not found. Make sure it's in the same directory.")
        return 0
    except Exception as e:
        print(f"Scan error: {e}")
        return 0

def process_files():
    """Process files using embedder"""
    try:
        from embedder import process_files
        print("Running file processing with embedder...")
        processed = process_files(limit=50)
        print(f"Processed {processed} files")
        return processed
    except ImportError:
        print("Error: embedder module not found.")
        return 0
    except Exception as e:
        print(f"Processing error: {e}")
        return 0

def get_database_stats():
    """Get database statistics"""
    try:
        conn = get_db_conn()
        cursor = conn.cursor()
        
        stats = {}
        
        # File counts
        cursor.execute("SELECT COUNT(*) FROM files")
        stats['total_files'] = cursor.fetchone()[0]
        
        cursor.execute("SELECT COUNT(*) FROM files WHERE processed = 1")
        stats['processed_files'] = cursor.fetchone()[0]
        
        cursor.execute("SELECT COUNT(*) FROM files WHERE processed = -1")
        stats['failed_files'] = cursor.fetchone()[0]
        
        cursor.execute("SELECT COUNT(*) FROM files WHERE processed = 0")
        stats['unprocessed_files'] = cursor.fetchone()[0]
        
        # File type breakdown
        cursor.execute("SELECT type, COUNT(*) FROM files GROUP BY type")
        stats['file_types'] = dict(cursor.fetchall())
        
        # Vector counts by modality
        cursor.execute("SELECT modality, COUNT(*) FROM vectors GROUP BY modality")
        stats['vectors'] = dict(cursor.fetchall())
        
        # Total vectors
        cursor.execute("SELECT COUNT(*) FROM vectors")
        stats['total_vectors'] = cursor.fetchone()[0]
        
        conn.close()
        return stats
        
    except Exception as e:
        return {"error": str(e)}

def show_status():
    """Display system status"""
    stats = get_database_stats()
    
    if "error" in stats:
        print(f"Error getting stats: {stats['error']}")
        return
    
    print("=" * 60)
    print("IMPROVED FILE CHATBOT STATUS - LLAVA RESTRICTED MODE")
    print("=" * 60)
    
    # Configuration
    print(f"Configuration:")
    print(f"  Default top-k results: {DEFAULT_TOP_K}")
    print(f"  Max top-k results: {MAX_TOP_K}")
    print(f"  Memory turns: {MEMORY_TURNS}")
    print(f"  Text model: {TEXT_MODEL}")
    print(f"  Multi-modal model: {MULTI_MODEL} (RESTRICTED to standalone images/videos)")
    print(f"  Llava usage: Only for standalone image/video files")
    print(f"  Text files: Processed with {TEXT_MODEL} + OCR only")
    
    # Files
    print(f"\nFiles in Database: {stats['total_files']}")
    print(f"  Processed: {stats['processed_files']}")
    print(f"  Failed: {stats['failed_files']}")
    print(f"  Unprocessed: {stats['unprocessed_files']}")
    
    # File types
    if stats['file_types']:
        print(f"\nFile Types:")
        for ftype, count in sorted(stats['file_types'].items()):
            if ftype:
                processing_method = "Llava" if ftype in ['image', 'video'] else "Text Model + OCR"
                print(f"  {ftype}: {count} files → {processing_method}")
    
    # Vector embeddings
    print(f"\nVector Embeddings: {stats['total_vectors']}")
    if stats['vectors']:
        for modality, count in sorted(stats['vectors'].items()):
            print(f"  {modality}: {count}")
    
    # FAISS indices
    print(f"\nFAISS Indices:")
    text_idx = Path("indices/faiss_text.index")
    image_idx = Path("indices/faiss_image.index")
    print(f"  Text index: {'Present' if text_idx.exists() else 'Missing'}")
    print(f"  Image index: {'Present' if image_idx.exists() else 'Missing'}")
    
    print("=" * 60)

def initialize_database():
    """Initialize database tables if they don't exist"""
    conn = get_db_conn()
    cursor = conn.cursor()
    
    # Create files table (compatible with embedder)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS files (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            path TEXT UNIQUE NOT NULL,
            type TEXT,
            processed INTEGER DEFAULT 0,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        )
    """)
    
    # Create vectors table (compatible with embedder)
    cursor.execute("""
        CREATE TABLE IF NOT EXISTS vectors (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            file_id INTEGER,
            chunk_id INTEGER DEFAULT 0,
            modality TEXT NOT NULL,
            embedding BLOB,
            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
            FOREIGN KEY (file_id) REFERENCES files (id)
        )
    """)
    
    # Create indices for performance
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_files_path ON files(path)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_files_processed ON files(processed)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_vectors_file_id ON vectors(file_id)")
    cursor.execute("CREATE INDEX IF NOT EXISTS idx_vectors_modality ON vectors(modality)")
    
    conn.commit()
    conn.close()

# Main function
def main():
    """Main function"""
    print("IMPROVED FILE CHATBOT - LLAVA RESTRICTED MODE")
    print("Llava only processes standalone image/video files")
    print("Text files (PDF, DOCX, TXT) use text model + OCR only")
    print("=" * 70)
    
    # Initialize database
    initialize_database()
    
    print(f"\nCommands:")
    print(f"- 'scan /path/to/directory' - Add files to database")
    print(f"- 'find keyword' or 'search keyword' - Search files")
    print(f"- 'find videos top 3' - Search with custom result limit")
    print(f"- 'find badminton videos' - Enhanced search (filename + semantic)")
    print(f"- 'read N' or 'analyze N' - Process file number N")
    print(f"  • Images/Videos: Llava multimodal analysis")
    print(f"  • Documents: Text model + OCR analysis")
    print(f"- 'read-ocr N' - Show raw OCR output for file N")
    print(f"- 'extract data from 1' - Extract specific information")
    print(f"- 'process' - Process unindexed files with embedder")
    print(f"- 'status' - Show system status")
    print(f"- 'top 3' or 'top 8' - Set number of results (max {MAX_TOP_K})")
    print(f"- 'quit' or 'exit' - Exit")
    print("=" * 70)
    
    # Show initial status
    show_status()
    
    # Check for indices
    if not Path("indices").exists():
        print(f"\nWarning: indices/ directory not found")
        print(f"FAISS search may not work until you run 'process' to create indices")

    print(f"\nSystem ready! Default: top {DEFAULT_TOP_K} results")
    print("Llava restricted to standalone image/video files only")
    
    # Initialize improved chatbot
    bot = ImprovedFileChatbot()

    # Main chat loop
    while True:
        try:
            user_input = input(f"\nYou: ").strip()
            if not user_input:
                continue
            
            if user_input.lower() in ('quit', 'exit', 'q'):
                print("Goodbye!")
                break
            
            # Handle special commands
            if user_input.lower().startswith('scan '):
                directory = user_input[5:].strip().strip('"\'')
                scan_directory(directory)
                continue
            
            elif user_input.lower() == 'process':
                process_files()
                continue
            
            elif user_input.lower() == 'status':
                show_status()
                continue
            
            elif user_input.lower().startswith('top '):
                try:
                    k = int(user_input[4:].strip())
                    bot.set_top_k(k)
                    print(f"Set to return top {bot.top_k} results")
                    continue
                except ValueError:
                    print("Invalid number. Use 'top N' where N is a number.")
                    continue
            
            elif user_input.lower() in ['help', '?']:
                print(f"\nAvailable commands:")
                print(f"- scan /path - Add files from directory")
                print(f"- find/search keyword - Search files")
                print(f"- find videos top 3 - Search with custom limit")
                print(f"- read/analyze N - Process file number N")
                print(f"- extract data from N - Extract specific info")
                print(f"- read-ocr N - Show raw OCR output")
                print(f"- top N - Set result limit (max {MAX_TOP_K})")
                continue
            
            # Process regular chat/search
            response = bot.process(user_input)
            print(f"\nBot: {response}")
            
        except KeyboardInterrupt:
            print(f"\nGoodbye!")
            break
        except Exception as e:
            print(f"Error: {e}")
            import traceback
            traceback.print_exc()

if __name__ == "__main__":
    main()